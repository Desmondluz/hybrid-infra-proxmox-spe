"""
Build script for the CIA — Hybrid Infrastructure keynote pptx.

Reproducible build of the final keynote presentation for the FW3 Final.
Output: docs/keynote/CIA-GR46-keynote-final.pptx

Palette : Midnight Executive
  - navy       #1E2761  (dominant, dark slides)
  - ice blue   #CADCFC  (light bg + accents)
  - white      #FFFFFF
  - accent     #38BDF8  (CIA project accent, from keynote.md)
  - success    #22C55E
  - alert      #EF4444

Usage:
    python3 scripts/keynote/build.py

Requires: pip install python-pptx Pillow
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

REPO = Path(__file__).resolve().parents[2]
CAPTURES = REPO / "docs" / "demo" / "captures"
OUT_DIR = REPO / "docs" / "keynote"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT = OUT_DIR / "CIA-GR46-keynote-final.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette Midnight Executive
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)  # CIA bright blue
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)
ALERT = RGBColor(0xEF, 0x44, 0x44)
SLATE = RGBColor(0x64, 0x74, 0x8B)  # secondary text
DARK_NAVY = RGBColor(0x10, 0x16, 0x3A)  # darker accent

# Fonts
TITLE_FONT = "Trebuchet MS"
BODY_FONT = "Calibri"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    left, top, width, height,
    *,
    font_size=14,
    bold=False,
    italic=False,
    color: RGBColor = NAVY,
    font_name: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """Add a text box at given position with formatted text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multi_text(
    slide,
    paragraphs,  # list of dicts: {text, size, bold, color, align, font}
    left, top, width, height,
    *,
    anchor=MSO_ANCHOR.TOP,
):
    """Add a text box with multiple paragraphs of varying styles."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = para["text"]
        run.font.name = para.get("font", BODY_FONT)
        run.font.size = Pt(para.get("size", 14))
        run.font.bold = para.get("bold", False)
        run.font.italic = para.get("italic", False)
        run.font.color.rgb = para.get("color", NAVY)
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
    return tb


def add_rect(slide, left, top, width, height, fill: RGBColor, line=None, rounded=False):
    """Add a rectangle shape (filled, optional line color)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    return rect


def add_image(slide, image_path: Path, left, top, width=None, height=None):
    """Add an image to the slide."""
    if not image_path.exists():
        print(f"  ! Missing image: {image_path}")
        return None
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(image_path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(image_path), left, top, height=height)
    return slide.shapes.add_picture(str(image_path), left, top)


def add_speaker_notes(slide, notes: str):
    """Add presenter notes to the slide."""
    nslide = slide.notes_slide
    tf = nslide.notes_text_frame
    tf.text = notes


def add_blank_slide(prs):
    """Create a new blank slide and return it."""
    blank_layout = prs.slide_layouts[6]  # layout 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_footer(slide, page_num: int, total: int = 15, dark=False):
    """Add a small footer with page numbers and project tag."""
    color = ICE_BLUE if dark else SLATE
    add_text(
        slide, "CIA · Hybrid Infrastructure · GR46",
        Inches(0.5), Inches(7.05), Inches(6), Inches(0.35),
        font_size=9, color=color, align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, f"{page_num} / {total}",
        Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35),
        font_size=9, color=color, align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Slide 1 — Title (dark sandwich top)."""
    s = add_blank_slide(prs)
    set_bg(s, NAVY)

    # Accent vertical bar on the left
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)

    # Top tag
    add_text(s, "EPITECH MSc Pro · 2025-2026 · T-NSA-810-REP25",
             Inches(0.7), Inches(0.5), Inches(10), Inches(0.4),
             font_size=12, color=ICE_BLUE, italic=True)

    # Main title
    add_text(s, "CIA",
             Inches(0.7), Inches(1.3), Inches(12), Inches(1.2),
             font_size=96, bold=True, color=WHITE, font_name=TITLE_FONT)

    add_text(s, "Hybrid Infrastructure",
             Inches(0.7), Inches(2.5), Inches(12), Inches(0.9),
             font_size=44, bold=True, color=ICE_BLUE, font_name=TITLE_FONT)

    # Subtitle
    add_text(s, "Deployment & Securing — Proxmox · pfSense · Elastic",
             Inches(0.7), Inches(3.45), Inches(12), Inches(0.5),
             font_size=22, color=ICE_BLUE)

    # Tagline
    add_text(s, "GitOps · Defense in Depth · Runtime Observability",
             Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
             font_size=16, italic=True, color=ACCENT)

    # Footer block - presenter
    add_rect(s, Inches(0.7), Inches(5.8), Inches(5.5), Inches(0.05), ACCENT)
    add_text(s, "Présenté par GR46",
             Inches(0.7), Inches(5.95), Inches(8), Inches(0.4),
             font_size=14, bold=True, color=WHITE)
    add_text(s, "Joseph-Desmon Yonzou · dyonzou@gmail.com",
             Inches(0.7), Inches(6.35), Inches(8), Inches(0.35),
             font_size=12, color=ICE_BLUE)
    add_text(s, "Final · Juin 2026",
             Inches(0.7), Inches(6.7), Inches(8), Inches(0.35),
             font_size=12, color=ICE_BLUE)

    # Right-side stat tags
    for i, (label, val) in enumerate([
        ("2", "sites"),
        ("6+1", "VMs"),
        ("4", "couches"),
        ("15K", "events"),
    ]):
        x = Inches(7.2 + i * 1.45)
        add_text(s, label, x, Inches(5.95), Inches(1.4), Inches(0.6),
                 font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)
        add_text(s, val, x, Inches(6.55), Inches(1.4), Inches(0.4),
                 font_size=11, color=ICE_BLUE, align=PP_ALIGN.CENTER)

    add_speaker_notes(s, (
        "Bonjour, je suis Joseph-Desmon Yonzou du groupe GR46. "
        "Je vous présente aujourd'hui le projet CIA — Hybrid Infrastructure, "
        "rendu T-NSA-810. L'objectif : déployer et sécuriser une infrastructure "
        "hybride sur 2 sites distants, avec une approche GitOps stricte (tout-code, "
        "tout-version, tout-audit), une défense en profondeur, et une observabilité "
        "runtime live. La présentation dure 20 minutes, suivies de 10 minutes "
        "de Q&R. À droite, les chiffres-clés du projet livré."
    ))


def slide_02_mandate(prs):
    """Slide 2 — Le mandat."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    # Title
    add_text(s, "Le mandat",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=40, bold=True, color=NAVY, font_name=TITLE_FONT)

    add_text(s, "Une entreprise déploie 2 sites distants — elle doit garantir...",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=18, italic=True, color=SLATE)

    # 3 pillar cards
    pillars = [
        ("C", "Confidentialité", "Chiffrement bout-à-bout, secrets gérés (SOPS + age), "
         "accès SSH bastion MFA, exception forward TCP nominative.", ACCENT),
        ("I", "Intégrité", "Configurations idempotentes (Terraform + Ansible), "
         "audit immutable des changes, killswitch d'isolation, secrets versionés chiffrés.", SUCCESS),
        ("A", "Audit", "Traçabilité totale : NetBox = source de vérité IPAM, "
         "logs centralisés Filebeat → Elastic, alerting actif, runbooks DRP.", RGBColor(0xF5, 0x9E, 0x0B)),
    ]

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    start_x = Inches(0.5)

    for i, (letter, title, desc, color) in enumerate(pillars):
        x = start_x + (card_w + gap) * i
        # Card background
        add_rect(s, x, Inches(2.0), card_w, card_h, ICE_BLUE, rounded=True)
        # Big letter circle (top of card)
        circle_size = Inches(1.4)
        cx = x + (card_w - circle_size) / 2
        add_rect(s, cx, Inches(2.3), circle_size, circle_size, color, rounded=True)
        add_text(s, letter, cx, Inches(2.35), circle_size, circle_size,
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name=TITLE_FONT)
        # Title
        add_text(s, title, x + Inches(0.25), Inches(4.0), card_w - Inches(0.5), Inches(0.6),
                 font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)
        # Description
        add_text(s, desc, x + Inches(0.3), Inches(4.65), card_w - Inches(0.6), Inches(1.7),
                 font_size=12, color=NAVY, align=PP_ALIGN.LEFT)

    # Bottom highlight
    add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), NAVY, rounded=True)
    add_text(s, "Parti pris GR46 : tout-code · tout-version (git) · tout-audit",
             Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.4),
             font_size=14, bold=True, color=ICE_BLUE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 2)

    add_speaker_notes(s, (
        "Le mandat client est classique : 2 sites distants, garantir le triptyque "
        "CIA — confidentialité, intégrité, audit. Notre parti pris est radical : "
        "tout est code (Terraform + Ansible), tout est versionné dans Git, et "
        "tout est auditable. Aucune action manuelle non reproductible. C'est ce qui "
        "fait la différence entre un projet 'qui marche' et un projet 'qui se reproduit'."
    ))


def slide_03_architecture(prs):
    """Slide 3 — Architecture globale (custom shapes)."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    # Title
    add_text(s, "Architecture globale",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=40, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "2 sites Proxmox · 3 VMs/site · tunnel OpenVPN site-à-site",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=16, italic=True, color=SLATE)

    # Site A box (left)
    site_w = Inches(5.5)
    site_h = Inches(3.8)
    siteA_x = Inches(0.5)
    siteA_y = Inches(2.0)

    add_rect(s, siteA_x, siteA_y, site_w, site_h, ICE_BLUE, rounded=True)
    add_text(s, "Site A · ns3050272.ip-51-255-76.eu",
             siteA_x + Inches(0.2), siteA_y + Inches(0.15), site_w - Inches(0.4), Inches(0.4),
             font_size=14, bold=True, color=NAVY)
    add_text(s, "Hub OpenVPN · PVE 9.1.4 · LAN 10.10.0.0/24",
             siteA_x + Inches(0.2), siteA_y + Inches(0.55), site_w - Inches(0.4), Inches(0.3),
             font_size=10, color=SLATE, italic=True)

    vm_w = Inches(1.55)
    vm_h = Inches(1.0)
    vm_y = siteA_y + Inches(1.0)
    for i, (name, role) in enumerate([
        ("pfsense-s1", "FW · VPN server"),
        ("bastion-s1", "SSH MFA"),
        ("services-s1", "NetBox · Vault"),
    ]):
        x = siteA_x + Inches(0.2) + (vm_w + Inches(0.1)) * i
        add_rect(s, x, vm_y, vm_w, vm_h, WHITE, rounded=True)
        add_text(s, name, x, vm_y + Inches(0.1), vm_w, Inches(0.35),
                 font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, role, x, vm_y + Inches(0.45), vm_w, Inches(0.4),
                 font_size=9, color=SLATE, align=PP_ALIGN.CENTER)

    # Obs runtime box for site A
    obs_y = siteA_y + Inches(2.3)
    add_rect(s, siteA_x + Inches(0.2), obs_y, site_w - Inches(0.4), Inches(1.3), WHITE, rounded=True)
    add_text(s, "Elastic stack · Logstash · Kibana 8.11", siteA_x + Inches(0.3), obs_y + Inches(0.15),
             site_w - Inches(0.5), Inches(0.3), font_size=11, bold=True, color=NAVY)
    add_text(s, "Dashboard runtime + 3 alert rules · pipeline Filebeat live",
             siteA_x + Inches(0.3), obs_y + Inches(0.5), site_w - Inches(0.5), Inches(0.7),
             font_size=10, color=SLATE)

    # Site B box (right)
    siteB_x = Inches(7.3)
    siteB_y = siteA_y
    add_rect(s, siteB_x, siteB_y, site_w, site_h, ICE_BLUE, rounded=True)
    add_text(s, "Site B · ns3183326.ip-146-59-253.eu",
             siteB_x + Inches(0.2), siteB_y + Inches(0.15), site_w - Inches(0.4), Inches(0.4),
             font_size=14, bold=True, color=NAVY)
    add_text(s, "Client OpenVPN · PVE 9.1.14 · LAN 192.168.0.0/24",
             siteB_x + Inches(0.2), siteB_y + Inches(0.55), site_w - Inches(0.4), Inches(0.3),
             font_size=10, color=SLATE, italic=True)

    for i, (name, role) in enumerate([
        ("pfsense-s2", "FW · VPN client"),
        ("bastion-s2", "SSH MFA"),
        ("services-s2", "Bastion-lite · Obs"),
    ]):
        x = siteB_x + Inches(0.2) + (vm_w + Inches(0.1)) * i
        add_rect(s, x, vm_y, vm_w, vm_h, WHITE, rounded=True)
        add_text(s, name, x, vm_y + Inches(0.1), vm_w, Inches(0.35),
                 font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, role, x, vm_y + Inches(0.45), vm_w, Inches(0.4),
                 font_size=9, color=SLATE, align=PP_ALIGN.CENTER)

    # Site C cloud
    add_rect(s, siteB_x + Inches(0.2), siteB_y + Inches(2.3), site_w - Inches(0.4), Inches(1.3),
             RGBColor(0xFE, 0xF3, 0xC7), rounded=True)
    add_text(s, "Site C cloud Azure · germanywestcentral (planifié)",
             siteB_x + Inches(0.3), siteB_y + Inches(2.45), site_w - Inches(0.5), Inches(0.3),
             font_size=11, bold=True, color=NAVY)
    add_text(s, "Code Terraform livré · 7/8 ressources runtime · VM bloquée par capacity Students",
             siteB_x + Inches(0.3), siteB_y + Inches(2.8), site_w - Inches(0.5), Inches(0.7),
             font_size=10, color=SLATE)

    # Tunnel VPN line between sites
    tunnel_y = Inches(3.7)
    tunnel_line = add_rect(s, Inches(6.0), tunnel_y, Inches(1.3), Inches(0.08), ACCENT)
    add_text(s, "OpenVPN AES-256-GCM",
             Inches(5.5), tunnel_y - Inches(0.35), Inches(2.3), Inches(0.3),
             font_size=10, italic=True, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, "tunnel site-à-site",
             Inches(5.5), tunnel_y + Inches(0.15), Inches(2.3), Inches(0.3),
             font_size=9, color=SLATE, align=PP_ALIGN.CENTER)

    # Footer
    add_footer(s, 3)

    add_speaker_notes(s, (
        "Voici notre architecture cible. Deux clusters Proxmox loués à OVH "
        "(ns3050272 = Site A, ns3183326 = Site B), trois VMs sur chaque site, "
        "reliés par un tunnel OpenVPN site-à-site chiffré en AES-256-GCM. "
        "Site A héberge le hub VPN, NetBox et Vault. Site B héberge nos services "
        "et fait office de site secondaire. La couche observabilité Elastic+Kibana "
        "tourne sur le poste de pilotage en WSL pour cette démo (via tunnel SSH "
        "inverse, on y reviendra). Le Site C Azure était notre bonus cloud — "
        "code 100% prêt, 7 ressources réseau créées en runtime, seule la VM "
        "Linux bloquée par une pénurie de capacité B-series Azure ce week-end."
    ))


def slide_04_stack(prs):
    """Slide 4 — Stack technique (2x3 grid)."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Stack technique",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=40, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "6 couches · 100 % open source · 100 % infra-as-code",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=16, italic=True, color=SLATE)

    # 2x3 grid of stack items
    items = [
        ("Terraform 1.10", "Provisioning",
         "modules Proxmox + Azure azurerm 3.117 · plan-then-apply · state local + roadmap remote backend"),
        ("Ansible 2.18", "Configuration",
         "11 rôles : common, pfsense, openvpn, netbox, elastic/kibana/logstash, bastion, dns, webapp, filebeat"),
        ("pfSense CE 2.7", "Firewall as code",
         "rules sourcées XML · DNAT bastion · killswitch floating rule · matrice 28 flux autorisés / 8 refus"),
        ("NetBox 4.0", "IPAM source de vérité",
         "sites + prefixes + IP allocations · automation Ansible · auto-sync depuis inventaire"),
        ("Elastic 8.11", "Observabilité runtime",
         "Filebeat agents · Logstash pipeline · ES single-node · Kibana dashboards + 3 alert rules live"),
        ("SOPS + age", "Secrets chiffrés",
         "aucun secret en clair Git · chiffrement par destinataire age · déchiffrement Ansible runtime"),
    ]

    cell_w = Inches(4.1)
    cell_h = Inches(2.55)
    start_x = Inches(0.5)
    start_y = Inches(2.0)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)

    for i, (name, role, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = start_x + (cell_w + gap_x) * col
        y = start_y + (cell_h + gap_y) * row

        add_rect(s, x, y, cell_w, cell_h, ICE_BLUE, rounded=True)

        # Accent strip on the left
        add_rect(s, x, y, Inches(0.15), cell_h, ACCENT)

        # Name (tool + version)
        add_text(s, name, x + Inches(0.35), y + Inches(0.15), cell_w - Inches(0.5), Inches(0.45),
                 font_size=18, bold=True, color=NAVY, font_name=TITLE_FONT)
        # Role tag
        add_text(s, role.upper(), x + Inches(0.35), y + Inches(0.65), cell_w - Inches(0.5), Inches(0.35),
                 font_size=11, bold=True, color=ACCENT)
        # Description
        add_text(s, desc, x + Inches(0.35), y + Inches(1.05), cell_w - Inches(0.55), cell_h - Inches(1.2),
                 font_size=11, color=NAVY)

    add_footer(s, 4)

    add_speaker_notes(s, (
        "La stack est volontairement classique et open source. Terraform pour le "
        "provisioning (avec un module Proxmox custom et un module Azure pour le "
        "bonus cloud). Ansible pour la configuration : 11 rôles couvrant tous "
        "les services. pfSense pour le firewall, configuration sourcée en XML "
        "et 28 flux autorisés vs 8 refus dans la matrice. NetBox comme source "
        "de vérité IPAM. Elastic 8.11 pour l'observabilité runtime — c'est la "
        "couche qu'on a livrée cette semaine et que je vais démontrer. "
        "Et SOPS + age pour le chiffrement des secrets : aucun mot de passe "
        "en clair dans Git, chiffrement par destinataire, déchiffrement Ansible "
        "à la volée."
    ))


def slide_05_addressing(prs):
    """Slide 5 — Plan d'adressage + matrice flux."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Adressage & matrice de flux",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "Zones réseau · 28 flux autorisés explicitement · 8 flux refusés",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Left column: zones
    add_text(s, "Zones réseau",
             Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.5),
             font_size=18, bold=True, color=NAVY)

    zones = [
        ("A-L", "Site A LAN", "10.10.0.0/24"),
        ("A-M", "Site A ADMIN", "10.10.10.0/24"),
        ("A-W", "Site A WAN", "public OVH"),
        ("B-L", "Site B LAN", "192.168.0.0/24"),
        ("B-S", "Site B SERVICES", "192.168.10.0/24"),
        ("B-W", "Site B WAN", "public OVH"),
        ("VPN", "Tunnel site-à-site", "172.16.0.0/30"),
        ("INT", "Internet public", "0.0.0.0/0"),
    ]

    row_h = Inches(0.4)
    for i, (code, name, net) in enumerate(zones):
        y = Inches(2.5) + row_h * i
        bg = ICE_BLUE if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.5), y, Inches(5.5), row_h, bg)
        add_text(s, code, Inches(0.65), y, Inches(0.8), row_h,
                 font_size=12, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, Inches(1.45), y, Inches(2.5), row_h,
                 font_size=11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, net, Inches(3.95), y, Inches(2.0), row_h,
                 font_size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    # Right column: example flows
    add_text(s, "Exemples de flux autorisés",
             Inches(6.5), Inches(2.0), Inches(6), Inches(0.5),
             font_size=18, bold=True, color=NAVY)

    flows = [
        ("INT → B-W:2222", "TCP", "SSH bastion via DNAT pfSense", SUCCESS),
        ("A-M → B-S", "TCP/22", "Admin Site A vers services Site B via VPN", SUCCESS),
        ("A-L → INT:443", "TCP", "HTTPS sortant (stateful)", SUCCESS),
        ("B-S → A-L:9200", "TCP", "Filebeat → Elasticsearch (encrypted)", SUCCESS),
        ("INT → A-L", "ANY", "Refusé sauf retours stateful + OpenVPN", ALERT),
        ("B-W → A-M", "ANY", "Refusé : routing isolé par zone", ALERT),
    ]

    for i, (route, proto, desc, color) in enumerate(flows):
        y = Inches(2.5) + row_h * i
        bg = ICE_BLUE if i % 2 == 0 else WHITE
        add_rect(s, Inches(6.5), y, Inches(6.3), row_h, bg)
        add_text(s, route, Inches(6.65), y, Inches(2.3), row_h,
                 font_size=11, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, proto, Inches(8.95), y, Inches(0.8), row_h,
                 font_size=10, color=color, anchor=MSO_ANCHOR.MIDDLE, bold=True)
        add_text(s, desc, Inches(9.75), y, Inches(3.0), row_h,
                 font_size=10, color=SLATE, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom note
    add_text(s,
             "→ docs/access-matrix-network-flows.md : matrice complète + justifications par flux",
             Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             font_size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(s, 5)

    add_speaker_notes(s, (
        "Le plan d'adressage est rigoureux. 8 zones réseau distinctes, chacune "
        "avec son CIDR et sa fonction. Le tunnel VPN site-à-site sur 172.16.0.0/30. "
        "À droite, quelques exemples de flux : on autorise explicitement ce qui est "
        "nécessaire (DNAT bastion, admin cross-site via VPN, sortie HTTPS), on "
        "refuse tout le reste par défaut. La matrice complète documente 28 flux "
        "autorisés et 8 refus explicites — chacun avec sa justification métier. "
        "C'est dans access-matrix-network-flows.md du repo."
    ))


def slide_06_terraform_demo(prs):
    """Slide 6 — Démo Terraform idempotence + récit dev → école."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Démo Terraform · idempotence",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "Plan-then-apply · convergence dev jetable → matériel école sans réécriture",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Left: image of terraform plan
    img = CAPTURES / "07-terraform-plan-noop.png"
    add_image(s, img, Inches(0.5), Inches(2.0), width=Inches(7.0))

    # Right: récit migration
    right_x = Inches(7.8)
    add_text(s, "Migration dev → école",
             right_x, Inches(2.0), Inches(5.2), Inches(0.5),
             font_size=20, bold=True, color=NAVY, font_name=TITLE_FONT)

    steps = [
        ("1", "Dev jetable", "Proxmox nested VMware\n6 VMs auto-créées par Terraform"),
        ("2", "Code validé", "Modules + rôles testés en CI\n12 checks checkov, terraform validate OK"),
        ("3", "Bascule école", "Modification : terraform.tfvars + inventaire + secrets\nLe code applicatif n'a pas bougé"),
        ("4", "Réconciliation", "6 VMs école pré-allouées importées via terraform import\nConvergence lisible, pas de big bang"),
    ]

    step_y = Inches(2.7)
    for i, (num, title, desc) in enumerate(steps):
        y = step_y + Inches(1.05) * i
        # Number circle
        add_rect(s, right_x, y, Inches(0.5), Inches(0.5), ACCENT, rounded=True)
        add_text(s, num, right_x, y + Inches(0.05), Inches(0.5), Inches(0.4),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)
        # Step title
        add_text(s, title, right_x + Inches(0.65), y, Inches(4.5), Inches(0.3),
                 font_size=13, bold=True, color=NAVY)
        # Step desc
        add_text(s, desc, right_x + Inches(0.65), y + Inches(0.3), Inches(4.5), Inches(0.7),
                 font_size=10, color=SLATE)

    add_footer(s, 6)

    add_speaker_notes(s, (
        "À gauche, la sortie de notre terraform plan en production : 'No changes. "
        "Your infrastructure matches the configuration.' C'est la signature de "
        "l'idempotence. À droite, le récit de notre migration dev vers école : "
        "on a d'abord validé tout le code sur un Proxmox nested (VMware sur un "
        "vieux laptop), puis on a basculé sur le matériel réel de l'école en "
        "ne changeant que trois fichiers : terraform.tfvars, l'inventaire Ansible, "
        "et les secrets SOPS. Les 6 VMs école étaient pré-allouées, on les a "
        "réconciliées via terraform import. Pas de big bang, une convergence "
        "lisible. C'est le pattern GitOps de migration."
    ))


def slide_07_ansible_demo(prs):
    """Slide 7 — Démo Ansible runtime + idempotence."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Démo Ansible · runtime + idempotence",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "11 rôles modulaires · check mode · preuve idempotence services-s2",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Left: Big idempotence stat
    box_x = Inches(0.5)
    box_y = Inches(2.0)
    add_rect(s, box_x, box_y, Inches(6.0), Inches(4.0), NAVY, rounded=True)
    # Label
    add_text(s, "Preuve idempotence", box_x, box_y + Inches(0.25),
             Inches(6.0), Inches(0.4),
             font_size=14, color=ICE_BLUE, align=PP_ALIGN.CENTER, italic=True)
    # Big number "4" centered
    add_text(s, "4", box_x, box_y + Inches(0.7),
             Inches(6.0), Inches(1.6),
             font_size=120, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font_name=TITLE_FONT)
    # Label "changed tasks" right under
    add_text(s, "changed tasks", box_x, box_y + Inches(2.35),
             Inches(6.0), Inches(0.4),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font_name=TITLE_FONT)
    # Subtitle "sur 41 OK"
    add_text(s, "sur 41 tâches OK · ~ 90 % stables au 2e run",
             box_x, box_y + Inches(2.8), Inches(6.0), Inches(0.4),
             font_size=13, color=ICE_BLUE, align=PP_ALIGN.CENTER, italic=True)
    # Bottom playrecap line
    add_text(s,
             "PLAY RECAP — services-s2 (Site B)",
             box_x, box_y + Inches(3.4), Inches(6.0), Inches(0.4),
             font_size=11, color=ICE_BLUE, align=PP_ALIGN.CENTER,
             font_name="Consolas")

    # Right: 11 roles
    right_x = Inches(7.0)
    add_text(s, "11 rôles Ansible",
             right_x, Inches(2.0), Inches(6), Inches(0.5),
             font_size=18, bold=True, color=NAVY)

    roles = [
        ("common", "Baseline OS + ntp + audit"),
        ("pfsense", "Firewall API config"),
        ("openvpn", "VPN serveur + clients"),
        ("netbox", "IPAM + auto-sync"),
        ("elasticsearch", "Cluster ES"),
        ("kibana", "Dashboards + alerts"),
        ("logstash", "Pipeline ingest"),
        ("bastion", "SSH MFA + sshd hardening"),
        ("dns-forwarder", "Unbound resolver"),
        ("webapp", "Caddy + app interne"),
        ("filebeat", "Log forwarder"),
    ]

    cell_y = Inches(2.5)
    cell_h = Inches(0.35)
    for i, (name, desc) in enumerate(roles):
        col = i // 6
        row = i % 6
        x = right_x + Inches(0.0) + col * Inches(3.0)
        y = cell_y + cell_h * row
        # Tiny accent dot
        add_rect(s, x, y + Inches(0.12), Inches(0.1), Inches(0.1), ACCENT, rounded=True)
        add_text(s, name, x + Inches(0.2), y, Inches(0.9), cell_h,
                 font_size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Consolas")
        add_text(s, desc, x + Inches(1.1), y, Inches(1.9), cell_h,
                 font_size=9, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 7)

    add_speaker_notes(s, (
        "À gauche, la preuve d'idempotence chiffrée. Sur services-s2, le second "
        "passage du playbook complet rapporte changed=4 sur 41 tâches OK. "
        "Soit environ 90 % de stabilité — les 4 changes restants sont des handlers "
        "type 'restart sshd' qu'on peut tracer un par un. À droite, les 11 rôles "
        "Ansible qui couvrent toute notre stack : baseline OS, firewall, VPN, IPAM, "
        "observabilité, bastion, DNS, webapp, et le forwarder Filebeat. Chaque rôle "
        "est versionné indépendamment, testable en CI, et appelable via tags."
    ))


def slide_08_security(prs):
    """Slide 8 — Sécurité defense in depth."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Sécurité · Defense in depth",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "4 couches superposées · aucun secret en clair · forward TCP nominal",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # 4 layers as a vertical onion-like stack
    layers = [
        ("L1", "Bastion SSH + MFA TOTP",
         "Point d'entrée unique · TOTP par utilisateur · jails fail2ban",
         ACCENT),
        ("L2", "sshd durci + Match user",
         "AllowTcpForwarding no globalement · exception nominative pour l'opérateur infra",
         RGBColor(0x06, 0xB6, 0xD4)),
        ("L3", "Secrets chiffrés (SOPS + age)",
         "Aucun mot de passe en clair Git · chiffrement par destinataire · déchiffrement runtime",
         RGBColor(0x8B, 0x5C, 0xF6)),
        ("L4", "Killswitch egress + audit",
         "Floating rule pfSense isolation egress · audit log tamper alerting · DRP testé",
         RGBColor(0xEC, 0x48, 0x99)),
    ]

    layer_h = Inches(1.1)
    layer_y = Inches(2.0)
    for i, (code, title, desc, color) in enumerate(layers):
        y = layer_y + layer_h * i + Inches(0.05) * i

        # Background card
        add_rect(s, Inches(0.5), y, Inches(12.3), layer_h, ICE_BLUE, rounded=True)
        # Color accent on the left
        add_rect(s, Inches(0.5), y, Inches(0.7), layer_h, color)
        # Code badge
        add_text(s, code, Inches(0.5), y, Inches(0.7), layer_h,
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name=TITLE_FONT)
        # Title
        add_text(s, title, Inches(1.4), y + Inches(0.15), Inches(10.8), Inches(0.45),
                 font_size=18, bold=True, color=NAVY, font_name=TITLE_FONT)
        # Description
        add_text(s, desc, Inches(1.4), y + Inches(0.55), Inches(10.8), Inches(0.5),
                 font_size=12, color=NAVY)

    add_footer(s, 8)

    add_speaker_notes(s, (
        "Notre approche sécurité est en défense en profondeur. Quatre couches qui "
        "se renforcent. Couche 1 : le bastion SSH avec MFA TOTP — point d'entrée "
        "unique, jails fail2ban. Couche 2 (nouvelle cette semaine) : sshd durci "
        "avec AllowTcpForwarding désactivé globalement, et une exception nominative "
        "Match User pour l'opérateur infra. C'est ce qui nous a permis de monter "
        "notre tunnel SSH inverse Filebeat tout en gardant le pivot SSH bloqué "
        "pour tous les autres utilisateurs. Couche 3 : SOPS + age. Aucun mot de "
        "passe en clair Git, chiffrement par destinataire. Couche 4 : killswitch "
        "egress (floating rule pfSense qui isole tout trafic sortant en un play "
        "Ansible) + audit log tamper alerting via Kibana."
    ))


def slide_09_killswitch(prs):
    """Slide 9 — Killswitch en action."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Killswitch · isolation egress en 1 play",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "Floating rule pfSense · activable / désactivable via Ansible · audit trail Git",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Left: Image of pfSense killswitch
    img = CAPTURES / "fw3-pfsense-siteB-floating-killswitch.png"
    add_image(s, img, Inches(0.5), Inches(2.0), width=Inches(7.5))

    # Right: command + flow
    right_x = Inches(8.3)
    add_text(s, "Activation en 1 commande",
             right_x, Inches(2.0), Inches(4.7), Inches(0.5),
             font_size=18, bold=True, color=NAVY, font_name=TITLE_FONT)

    # Command box
    add_rect(s, right_x, Inches(2.6), Inches(4.7), Inches(1.0), DARK_NAVY, rounded=True)
    add_text(s,
             "$ ansible-playbook killswitch.yml \\\n"
             "    -e killswitch_state=active \\\n"
             "    -e site=siteB",
             right_x + Inches(0.15), Inches(2.65), Inches(4.5), Inches(0.95),
             font_size=10, color=ACCENT, font_name="Consolas")

    # Effect
    add_text(s, "Effet immédiat", right_x, Inches(3.8), Inches(4.7), Inches(0.4),
             font_size=14, bold=True, color=NAVY)

    effects = [
        ("Avant", "curl → 200 OK", SUCCESS),
        ("Après", "curl → timeout", ALERT),
        ("Revert", "killswitch_state=inactive", ACCENT),
    ]
    for i, (label, val, color) in enumerate(effects):
        y = Inches(4.2) + Inches(0.5) * i
        add_text(s, label, right_x, y, Inches(1.3), Inches(0.4),
                 font_size=12, bold=True, color=NAVY)
        add_text(s, val, right_x + Inches(1.4), y, Inches(3.3), Inches(0.4),
                 font_size=11, color=color, font_name="Consolas")

    # Bottom usage note
    add_text(s,
             "Usage : isolation post-incident, exercice DRP, basculement traffic, conformité.",
             right_x, Inches(5.8), Inches(4.7), Inches(0.7),
             font_size=11, italic=True, color=SLATE)

    add_footer(s, 9)

    add_speaker_notes(s, (
        "Le killswitch est notre arme atomique en cas d'incident. À gauche, "
        "la règle floating pfSense qui bloque tout egress quand la variable "
        "KILLSWITCH_ACTIVE est positionnée. À droite, l'activation en une seule "
        "commande Ansible : on passe killswitch_state=active, le play modifie "
        "le pfsense.xml + reload firewall, et toute sortie est immédiatement "
        "bloquée. Avant : curl répond 200 OK. Après : timeout. Revert en passant "
        "killswitch_state=inactive. Usage type : isolation post-incident pour "
        "investigation forensique, exercice DRP, ou simple test de conformité."
    ))


def slide_10_obs_pipeline(prs):
    """Slide 10 — Observabilité Pipeline runtime ⭐."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Observabilité · pipeline runtime",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "Filebeat → tunnel SSH chiffré → Logstash → Elasticsearch → Kibana — LIVE",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=ACCENT, bold=True)

    # Pipeline diagram (boxes with arrows)
    pipeline = [
        ("services-s2", "Filebeat", "auth.log, syslog, audit", NAVY),
        ("Tunnel SSH", "port 5044", "chiffré, ProxyJump", ACCENT),
        ("WSL Docker", "Logstash 8.11", "filter hostname + tag", NAVY),
        ("Elastic", "Elasticsearch", "index cia-{hostname}-{date}", NAVY),
        ("Kibana 8.11", "Dashboards", "+ 3 alert rules live", ACCENT),
    ]
    box_w = Inches(2.3)
    box_h = Inches(1.4)
    box_y = Inches(2.2)
    start_x = Inches(0.4)
    arrow_w = Inches(0.2)
    spacing_x = box_w + arrow_w

    for i, (top, mid, bot, color) in enumerate(pipeline):
        x = start_x + spacing_x * i
        add_rect(s, x, box_y, box_w, box_h, color, rounded=True)
        add_text(s, top, x, box_y + Inches(0.2), box_w, Inches(0.35),
                 font_size=11, italic=True, color=ICE_BLUE, align=PP_ALIGN.CENTER)
        add_text(s, mid, x, box_y + Inches(0.55), box_w, Inches(0.4),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)
        add_text(s, bot, x, box_y + Inches(0.95), box_w, Inches(0.4),
                 font_size=9, color=ICE_BLUE, align=PP_ALIGN.CENTER)

        # Arrow between boxes (except after last)
        if i < len(pipeline) - 1:
            arrow_x = x + box_w + Inches(0.01)
            arrow_y = box_y + box_h / 2 - Inches(0.04)
            add_rect(s, arrow_x, arrow_y, arrow_w - Inches(0.02), Inches(0.08), ACCENT)

    # Big stat callouts below
    stats = [
        ("15 222", "events live indexés", "depuis services-s2 en 24h"),
        ("587", "rules executions", "alerting check toutes les minutes"),
        ("100 %", "open-source stack", "ES + Kibana + Logstash + Filebeat 8.11"),
    ]
    stat_y = Inches(4.2)
    stat_w = Inches(4.0)
    stat_h = Inches(2.0)
    for i, (val, label, sub) in enumerate(stats):
        x = Inches(0.5) + Inches(4.2) * i
        add_rect(s, x, stat_y, stat_w, stat_h, ICE_BLUE, rounded=True)
        add_text(s, val, x, stat_y + Inches(0.2), stat_w, Inches(1.0),
                 font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)
        add_text(s, label, x, stat_y + Inches(1.2), stat_w, Inches(0.4),
                 font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, sub, x, stat_y + Inches(1.6), stat_w, Inches(0.35),
                 font_size=10, italic=True, color=SLATE, align=PP_ALIGN.CENTER)

    add_footer(s, 10)

    add_speaker_notes(s, (
        "Voici LA partie différenciante du projet : la couche observabilité "
        "runtime. Le pipeline complet va de Filebeat sur services-s2, à travers "
        "un tunnel SSH inverse chiffré (parce que le tunnel VPN site-à-site "
        "n'était pas encore actif au moment de la démo, on a démontré le pattern "
        "avec un SSH reverse — sécurité identique), arrive sur Logstash dans "
        "Docker WSL, puis dans Elasticsearch single-node, et restitué dans "
        "Kibana 8.11. Les chiffres-clés : 15 222 events réels indexés depuis "
        "services-s2 sur 24h, 587 exécutions de rules alerting, 100 % open source. "
        "Le tunnel SSH inverse + le sshd durci avec exception nominative Match "
        "User, c'est exactement le pattern bastion durci en production."
    ))


def slide_11_dashboards_alerts(prs):
    """Slide 11 — Dashboards Kibana + Alerts."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Dashboards Kibana + Alert rules",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "2 dashboards versionnés .ndjson · 3 alert rules .ndjson · reproductible en 1 import",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Left: dashboard screenshot
    add_text(s, "Dashboard SSH Security Monitor",
             Inches(0.5), Inches(1.8), Inches(6.5), Inches(0.4),
             font_size=14, bold=True, color=NAVY)
    img = CAPTURES / "kibana-ssh-security-monitor.png"
    add_image(s, img, Inches(0.5), Inches(2.25), width=Inches(6.5))

    # Right: alert rules
    add_text(s, "3 alert rules opérationnelles",
             Inches(7.3), Inches(1.8), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=NAVY)
    img2 = CAPTURES / "kibana-3-alert-rules.png"
    add_image(s, img2, Inches(7.3), Inches(2.25), width=Inches(5.5))

    # Rules summary at bottom
    rules = [
        ("ALERT-SSH-FAIL-BURST", "P2", "count('Failed password') > 5 · 5 min · check 1 min", ACCENT),
        ("ALERT-AUDIT-TAMPER", "P1", "match /etc/ssh/ /etc/sudoers /etc/pam.d/ > 0 · 10 min", ALERT),
        ("ALERT-FILEBEAT-HEARTBEAT-LOST", "P2", "count(*) < 1 · 10 min · auto-déclenché 2x organique", RGBColor(0xF5, 0x9E, 0x0B)),
    ]
    rule_y = Inches(5.7)
    rule_h = Inches(0.4)
    for i, (name, sev, cond, color) in enumerate(rules):
        y = rule_y + rule_h * i
        add_rect(s, Inches(0.5), y, Inches(0.8), rule_h, color)
        add_text(s, sev, Inches(0.5), y, Inches(0.8), rule_h,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, Inches(1.4), y, Inches(4.5), rule_h,
                 font_size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Consolas")
        add_text(s, cond, Inches(6.0), y, Inches(6.8), rule_h,
                 font_size=10, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 11)

    add_speaker_notes(s, (
        "Deux dashboards Kibana versionnés en ndjson dans le repo, importables "
        "en une commande. À gauche, le dashboard SSH Security Monitor : 5 panneaux "
        "Lens — Failed passwords 24h, Invalid users 24h, Accepted successes, "
        "Timeline failures, et table forensique des derniers incidents. "
        "À droite, la liste des 3 alert rules actives. ALERT-SSH-FAIL-BURST en "
        "P2, déclenchée si plus de 5 Failed passwords en 5 minutes. ALERT-AUDIT-TAMPER "
        "en P1 critique, déclenchée si modification détectée sur les fichiers "
        "sensibles. Et ALERT-FILEBEAT-HEARTBEAT-LOST en P2, déclenchée si plus "
        "aucun event ne remonte depuis 10 minutes — celle-là s'est même déclenchée "
        "organiquement deux fois pendant les redémarrages Kibana, preuve qu'elle "
        "marche pour de vrai."
    ))


def slide_12_alerting_live(prs):
    """Slide 12 — Alerting LIVE prouvé ⭐⭐⭐."""
    s = add_blank_slide(prs)
    set_bg(s, NAVY)

    # Title
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_text(s, "ALERTING LIVE · preuve runtime",
             Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE, font_name=TITLE_FONT)
    add_text(s, "Burst injecté en démo · rule déclenchée · payload custom logué · 12 events détectés",
             Inches(0.7), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=ICE_BLUE)

    # Left: rule active screenshot
    add_text(s, "Rule ALERT-SSH-FAIL-BURST · onglet Alerts",
             Inches(0.5), Inches(1.85), Inches(6.5), Inches(0.4),
             font_size=12, bold=True, color=ICE_BLUE)
    img = CAPTURES / "kibana-rule-ssh-burst-active.png"
    add_image(s, img, Inches(0.5), Inches(2.3), width=Inches(6.4))

    # Right: server log payload
    add_text(s, "Server Log Kibana — payload custom",
             Inches(7.2), Inches(1.85), Inches(5.8), Inches(0.4),
             font_size=12, bold=True, color=ICE_BLUE)
    img2 = CAPTURES / "kibana-server-log-alerts-fired.png"
    add_image(s, img2, Inches(7.2), Inches(2.3), width=Inches(5.8))

    # Bottom: the log message itself, framed
    add_rect(s, Inches(0.5), Inches(6.05), Inches(12.5), Inches(0.95), DARK_NAVY, rounded=True)
    add_text(s,
             '🚨 ALERT-SSH-FAIL-BURST (P2) — 12 SSH Failed passwords detectees en 5 min sur infra CIA.',
             Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.4),
             font_size=12, color=ACCENT, font_name="Consolas", bold=True)
    add_text(s,
             'Date 2026-06-21T09:10:19.309Z. Investigation: Discover → message: "Failed password" → Last 15 min.',
             Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4),
             font_size=11, color=ICE_BLUE, font_name="Consolas")

    add_footer(s, 12, dark=True)

    add_speaker_notes(s, (
        "Pour vous prouver que l'alerting tourne en vrai, j'ai injecté 12 entrées "
        "fictives Failed password via logger sur services-s2 — IPs publiques de "
        "test (198.51.100.0/24 et 203.0.113.0/24, ranges réservés RFC5737, "
        "juridiquement neutres). En moins d'une minute, la rule a fait son check, "
        "détecté que count > 5 sur la fenêtre 5 minutes, est passée en status "
        "Active, et a déclenché l'action CIA Server Log. À gauche : la rule "
        "avec son alerte Active. À droite : le terminal qui montre le payload "
        "custom logué dans Kibana. En bas, la ligne brute capturée : "
        "'ALERT-SSH-FAIL-BURST P2, 12 Failed passwords detectees en 5 min, "
        "timestamp 2026-06-21T09:10:19.309Z'. C'est la preuve la plus forte du "
        "projet : un système d'alerting de prod qui se déclenche en live avec "
        "payload custom. En prod, on remplace le Server log connector par "
        "Slack ou Webhook, c'est une ligne de config à changer."
    ))


def slide_13_ci_quality(prs):
    """Slide 13 — Qualité & CI/CD."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Qualité & CI/CD",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "Pre-commit local · 5 workflows GitHub Actions · 163 runs · 100 % Phase A green",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # Image: CI green
    img = CAPTURES / "12-ci-actions-green.png"
    add_image(s, img, Inches(0.5), Inches(2.0), width=Inches(7.5))

    # Right: workflows + checks
    right_x = Inches(8.3)
    add_text(s, "5 workflows actifs", right_x, Inches(2.0), Inches(4.7), Inches(0.5),
             font_size=18, bold=True, color=NAVY, font_name=TITLE_FONT)

    workflows = [
        ("quality", "markdownlint · yamllint · terraform fmt"),
        ("security-scan", "checkov · gitleaks · trivy · tflint"),
        ("ansible", "ansible-lint · syntax-check · molecule (planifié)"),
        ("terraform", "terraform validate · plan dry-run · checkov"),
        ("release", "tag-driven · auto-changelog · pptx export"),
    ]
    wf_y = Inches(2.6)
    wf_h = Inches(0.55)
    for i, (name, checks) in enumerate(workflows):
        y = wf_y + wf_h * i
        add_rect(s, right_x, y, Inches(4.7), wf_h - Inches(0.05), ICE_BLUE, rounded=True)
        # Green check
        add_rect(s, right_x + Inches(0.1), y + Inches(0.1), Inches(0.3), Inches(0.3),
                 SUCCESS, rounded=True)
        add_text(s, "✓", right_x + Inches(0.1), y + Inches(0.07),
                 Inches(0.3), Inches(0.3),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, right_x + Inches(0.55), y + Inches(0.05),
                 Inches(1.5), Inches(0.25),
                 font_size=12, bold=True, color=NAVY,
                 font_name="Consolas")
        add_text(s, checks, right_x + Inches(0.55), y + Inches(0.27),
                 Inches(4.0), Inches(0.25),
                 font_size=9, color=SLATE)

    # Conventional commits + PR template note
    add_rect(s, right_x, Inches(5.5), Inches(4.7), Inches(1.0), NAVY, rounded=True)
    add_text(s, "Conventional Commits",
             right_x + Inches(0.15), Inches(5.6), Inches(4.4), Inches(0.3),
             font_size=12, bold=True, color=ICE_BLUE)
    add_text(s, "feat / fix / docs / chore + scope obligatoire\n"
                "PR template · reviewers · CODEOWNERS",
             right_x + Inches(0.15), Inches(5.9), Inches(4.4), Inches(0.6),
             font_size=10, color=ICE_BLUE)

    add_footer(s, 13)

    add_speaker_notes(s, (
        "Côté qualité, on est rigoureux. Pre-commit en local avec markdownlint, "
        "yamllint, terraform fmt, gitleaks. Cinq workflows GitHub Actions actifs : "
        "quality, security-scan, ansible, terraform, et release. 163 runs au "
        "compteur. Sur la Phase A — Observability livrée cette semaine — les 4 "
        "commits ont passé en CI verte. Le workflow security-scan utilise checkov "
        "pour Terraform, gitleaks pour les secrets, et trivy pour les images "
        "container. Tous nos commits suivent Conventional Commits : feat, fix, "
        "docs, chore, avec scope obligatoire. Et un PR template + CODEOWNERS "
        "pour structurer les revues."
    ))


def slide_14_golden_path(prs):
    """Slide 14 — Golden path nouveau site (on-prem vs cloud Azure)."""
    s = add_blank_slide(prs)
    set_bg(s, WHITE)

    add_text(s, "Golden path · nouveau site",
             Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
             font_size=36, bold=True, color=NAVY, font_name=TITLE_FONT)
    add_text(s, "On-prem Proxmox ou cloud Azure · matrice de décision · reproduction en 1 run",
             Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=SLATE)

    # 2 columns: on-prem vs cloud
    onprem_x = Inches(0.5)
    cloud_x = Inches(6.9)
    col_w = Inches(6.0)
    col_h = Inches(4.5)
    col_y = Inches(2.0)

    # On-prem column
    add_rect(s, onprem_x, col_y, col_w, col_h, NAVY, rounded=True)
    add_text(s, "ON-PREM PROXMOX",
             onprem_x, col_y + Inches(0.2), col_w, Inches(0.4),
             font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "École · convention · gratuit",
             onprem_x, col_y + Inches(0.65), col_w, Inches(0.3),
             font_size=12, italic=True, color=ICE_BLUE, align=PP_ALIGN.CENTER)

    onprem_pts = [
        ("⏱  1 jour", "convention école → provision"),
        ("💰 Inclus", "matériel école dédié"),
        ("⚡ < 5 ms", "latence LAN école"),
        ("📦 Templates", "Ubuntu/pfSense locaux"),
        ("✅ Stable", "services persistants critiques"),
    ]
    for i, (head, sub) in enumerate(onprem_pts):
        y = col_y + Inches(1.1) + Inches(0.6) * i
        add_text(s, head, onprem_x + Inches(0.4), y, col_w - Inches(0.8), Inches(0.3),
                 font_size=14, bold=True, color=WHITE)
        add_text(s, sub, onprem_x + Inches(0.4), y + Inches(0.3), col_w - Inches(0.8), Inches(0.25),
                 font_size=10, color=ICE_BLUE, italic=True)

    # Cloud column
    add_rect(s, cloud_x, col_y, col_w, col_h, DARK_NAVY, rounded=True)
    add_text(s, "CLOUD AZURE (Site C)",
             cloud_x, col_y + Inches(0.2), col_w, Inches(0.4),
             font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "Students subscription · code complet livré",
             cloud_x, col_y + Inches(0.65), col_w, Inches(0.3),
             font_size=12, italic=True, color=ICE_BLUE, align=PP_ALIGN.CENTER)

    cloud_pts = [
        ("⏱  5 min", "code Terraform azurerm prêt"),
        ("💰 ~10 €/mois", "B2s H24 · démo = quelques €"),
        ("⚡ 15-30 ms", "francecentral / germanywestcentral"),
        ("📦 Cloud-init", "Docker CE auto-installé"),
        ("⚠ 7/8 runtime", "VM bloquée capacity Students"),
    ]
    for i, (head, sub) in enumerate(cloud_pts):
        y = col_y + Inches(1.1) + Inches(0.6) * i
        add_text(s, head, cloud_x + Inches(0.4), y, col_w - Inches(0.8), Inches(0.3),
                 font_size=14, bold=True, color=WHITE)
        add_text(s, sub, cloud_x + Inches(0.4), y + Inches(0.3), col_w - Inches(0.8), Inches(0.25),
                 font_size=10, color=ICE_BLUE, italic=True)

    # Bottom highlight
    add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), ACCENT, rounded=True)
    add_text(s,
             "Règle GR46 : on-prem pour les services critiques · cloud off-site pour bonus DR + observabilité",
             Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.4),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 14)

    add_speaker_notes(s, (
        "Le golden path 'nouveau site' est documenté dans onboarding-new-site.md "
        "avec une matrice de décision on-prem vs cloud. À gauche, le pattern "
        "on-prem Proxmox : gratuit (convention école), latence sub-5ms, templates "
        "locaux, idéal pour les services critiques persistants. À droite, le "
        "pattern cloud Azure : code 100% prêt, déploiement en 5 minutes, coût "
        "~10€/mois pour une VM H24 — quelques euros sur la durée de démo. "
        "Note honnête : on a livré 7 ressources sur 8 en runtime Azure ce "
        "week-end, la VM finale est bloquée par une pénurie globale de capacité "
        "B-series Azure Europe. C'est un problème connu Microsoft, documenté, "
        "fixable en 24-72h via quota request. Règle GR46 : on-prem pour les "
        "services critiques de l'infra, cloud off-site pour le PRA et "
        "l'observabilité externe."
    ))


def slide_15_conclusion(prs):
    """Slide 15 — Bilan + ouverture (dark sandwich bottom)."""
    s = add_blank_slide(prs)
    set_bg(s, NAVY)

    # Accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)

    add_text(s, "Bilan & ouverture",
             Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             font_size=40, bold=True, color=WHITE, font_name=TITLE_FONT)
    add_text(s, "Ce qui marche · ce qui a coincé · ce qu'on ferait différemment",
             Inches(0.7), Inches(1.15), Inches(12), Inches(0.5),
             font_size=14, italic=True, color=ICE_BLUE)

    # 3 columns
    col_w = Inches(4.0)
    gap = Inches(0.2)
    col_y = Inches(2.0)
    col_h = Inches(4.5)

    cols = [
        ("WINS", SUCCESS, [
            ("GitOps marche",
             "Migration dev nested → école sans réécriture, juste 3 fichiers (tfvars, inv, secrets)"),
            ("Obs runtime live",
             "Pipeline Filebeat → tunnel SSH → Logstash → ES → Kibana, 15k+ events réels"),
            ("Alerting prouvé",
             "ALERT-SSH-FAIL-BURST déclenchée live, ALERT-FILEBEAT-HEARTBEAT déclenchée organiquement"),
            ("Defense in depth",
             "Bastion MFA + sshd Match user + SOPS+age + killswitch — 4 couches étanches"),
        ]),
        ("COINCES", RGBColor(0xF5, 0x9E, 0x0B), [
            ("Nested virt VMware",
             "VT-x à débloquer côté BIOS + Hyper-V → 2 jours perdus"),
            ("API pfSense",
             "courbe d'apprentissage (auth, XML structure, reload semantics)"),
            ("Disk full services-s2",
             "qm resize bloqué par permissions école → bastion en dual-duty"),
            ("Capacity Azure Students",
             "8 SKU × 3 régions testés, pénurie B-series globale ce week-end"),
        ]),
        ("FUTURE WORK", ACCENT, [
            ("Multi-site horizontal",
             "Site C-D en cascade via golden path, quota Azure obtenu"),
            ("OpenTofu migration",
             "Évaluer le fork open-source vs Terraform BSL"),
            ("Behavioral detection",
             "ML jobs Elastic + corrélation cross-host (bonus v3)"),
            ("Slack/Webhook connector",
             "Remplacer Server log par notif externe en prod"),
        ]),
    ]

    for i, (label, color, items) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        # Card
        add_rect(s, x, col_y, col_w, col_h, DARK_NAVY, rounded=True)
        # Color strip + label
        add_rect(s, x, col_y, col_w, Inches(0.6), color, rounded=True)
        # Re-draw the bottom edge of the strip squared (visual trick: another small rect)
        add_text(s, label, x, col_y + Inches(0.1), col_w, Inches(0.4),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font_name=TITLE_FONT)

        for j, (head, sub) in enumerate(items):
            iy = col_y + Inches(0.9) + Inches(0.9) * j
            add_text(s, head, x + Inches(0.3), iy, col_w - Inches(0.5), Inches(0.3),
                     font_size=12, bold=True, color=ACCENT)
            add_text(s, sub, x + Inches(0.3), iy + Inches(0.3), col_w - Inches(0.5), Inches(0.55),
                     font_size=9, color=ICE_BLUE, italic=True)

    # Bottom contact line
    add_rect(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.05), ACCENT)
    add_text(s, "Merci · Q&R 10 min",
             Inches(0.7), Inches(6.95), Inches(6), Inches(0.4),
             font_size=14, bold=True, color=WHITE)
    add_text(s, "github.com/Desmondluz/hybrid-infra-proxmox-spe",
             Inches(7), Inches(6.95), Inches(5.6), Inches(0.4),
             font_size=11, color=ICE_BLUE, align=PP_ALIGN.RIGHT,
             font_name="Consolas")

    add_speaker_notes(s, (
        "Pour conclure, le bilan honnête. Ce qui marche : la méthode GitOps "
        "a permis une migration dev jetable vers école avec juste 3 fichiers "
        "à changer. L'observabilité runtime est live, prouvée par 15 000 events "
        "réels et un alerting qui se déclenche pour de vrai. La défense en "
        "profondeur est étanche sur ses 4 couches. Ce qui a coincé : la "
        "nested virt VMware au début (2 jours perdus), l'API pfSense est "
        "subtile, le disk full sur services-s2 m'a forcé à un bastion en "
        "dual-duty, et la pénurie de capacité Azure B-series ce week-end "
        "a bloqué la VM Site C. Ce qu'on ferait différemment : figer les "
        "secrets SOPS dès J+1, écrire les runbooks en parallèle du code, "
        "demander le quota Azure dès le départ. Futur : Site C/D cloud en "
        "cascade, migration OpenTofu, détection comportementale, et connector "
        "Slack en prod. Merci, je prends vos questions."
    ))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_mandate,
        slide_03_architecture,
        slide_04_stack,
        slide_05_addressing,
        slide_06_terraform_demo,
        slide_07_ansible_demo,
        slide_08_security,
        slide_09_killswitch,
        slide_10_obs_pipeline,
        slide_11_dashboards_alerts,
        slide_12_alerting_live,
        slide_13_ci_quality,
        slide_14_golden_path,
        slide_15_conclusion,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  [{i:2d}/15] {builder.__name__}")
        builder(prs)

    prs.save(OUTPUT)
    print(f"\n  ✓ Saved: {OUTPUT}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size  : {SLIDE_W.inches}x{SLIDE_H.inches} (16:9)")


if __name__ == "__main__":
    build()
